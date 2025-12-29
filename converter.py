import io
import os
import tempfile
import fitz  # PyMuPDF
from docx import Document
from pdf2docx import Converter
import platform
import subprocess
import shutil
import time
from PIL import Image

def pdf_to_word(pdf_bytes):
    """
    Convert PDF bytes to Word document using pdf2docx for better layout preservation
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    
    print(f"DEBUG: Starting PDF to Word conversion... Size: {len(pdf_bytes)} bytes")
    temp_pdf = None
    temp_docx = None
    
    try:
        # Create temp files
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f:
            f.write(pdf_bytes)
            temp_pdf = f.name
            
        temp_docx = temp_pdf.replace(".pdf", ".docx")
        
        # Check if PDF is encrypted
        doc = fitz.open(temp_pdf)
        is_encrypted = doc.is_encrypted
        doc.close()
        
        if is_encrypted:
            raise Exception("This PDF is password protected. Please remove protection before converting.")

        # Convert using pdf2docx
        cv = Converter(temp_pdf)
        cv.convert(temp_docx, multi_processing=False)
        cv.close()
        
        # Read the generated DOCX
        if not os.path.exists(temp_docx):
            raise Exception("Conversion failed: Output file not created")
            
        print("DEBUG: pdf2docx conversion complete. Reading output...")
        with open(temp_docx, "rb") as f:
            docx_bytes = f.read()
            
        print("DEBUG: PDF to Word conversion successful.")
        return io.BytesIO(docx_bytes)
        
    except Exception as e:
        raise Exception(f"PDF to Word conversion failed: {str(e)}")
        
    finally:
        # Cleanup temp files
        safe_remove(temp_pdf)
        safe_remove(temp_docx)


def word_to_pdf_linux(docx_bytes):
    """
    Convert Word to PDF using LibreOffice (Headless) on Linux/Docker
    """
    temp_dir = tempfile.mkdtemp()
    temp_docx_path = os.path.join(temp_dir, "input.docx")
    
    try:
        with open(temp_docx_path, "wb") as f:
            f.write(docx_bytes)
            
        office_cmd = 'libreoffice'
        if shutil.which('soffice'):
            office_cmd = 'soffice'
        elif shutil.which('lowriter'):
            office_cmd = 'lowriter'

        cmd = [
            office_cmd, 
            '--headless', 
            '--convert-to', 
            'pdf', 
            '--outdir', 
            temp_dir, 
            temp_docx_path
        ]
        
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        
        temp_pdf_path = temp_docx_path.replace(".docx", ".pdf")
        
        if not os.path.exists(temp_pdf_path):
             error_msg = result.stderr.decode() if result.stderr else "Unknown error"
             raise Exception(f"LibreOffice conversion failed. Error: {error_msg}")
            
        with open(temp_pdf_path, "rb") as f:
            pdf_bytes = f.read()
            
        return io.BytesIO(pdf_bytes)
        
    except Exception as e:
        raise Exception(f"Linux Word to PDF conversion failed: {str(e)}")
    finally:
        if os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except:
                pass

def word_to_pdf(docx_bytes):
    """
    Convert Word to PDF using best available method:
    - Windows: Uses Microsoft Word (docx2pdf) for 100% fidelity.
    - Linux (Docker/Render): Uses LibreOffice for high compatibility.
    """
    print(f"DEBUG: Starting Word to PDF conversion... Size: {len(docx_bytes)} bytes")
    system = platform.system().lower()
    
    res = None
    if system == 'windows':
        res = word_to_pdf_windows(docx_bytes)
    else:
        res = word_to_pdf_linux(docx_bytes)
    
    print("DEBUG: Word to PDF conversion complete.")
    return res

def word_to_pdf_windows(docx_bytes):
    """
    Convert Word to PDF using Microsoft Word COM automation with robustness.
    """
    import pythoncom
    import win32com.client
    
    pythoncom.CoInitialize()
    temp_docx_path = None
    temp_pdf_path = None
    word = None
    doc = None
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_docx:
            temp_docx.write(docx_bytes)
            temp_docx_path = os.path.abspath(temp_docx.name)
            
        temp_pdf_path = temp_docx_path.replace(".docx", ".pdf")
        
        # Initialize Word with retry/robustness
        word = com_retry(win32com.client.DispatchEx, "Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0 # wdAlertsNone
        
        # Open document (ReadOnly=True, ConfirmConversions=False)
        doc = com_retry(word.Documents.Open, temp_docx_path, False, True)
        
        # Export as PDF (wdExportFormatPDF = 17)
        com_retry(doc.ExportAsFixedFormat, temp_pdf_path, 17)
        
        # Close and Cleanup
        com_retry(doc.Close, False)
        # We don't necessarily want to Quit Word if it was already open, 
        # but DispatchEx creates a new instance, so we should Quit it.
        com_retry(word.Quit)
        
        # Explicitly release handles
        del doc
        del word
        doc = None
        word = None
        
        if not os.path.exists(temp_pdf_path) or os.path.getsize(temp_pdf_path) == 0:
            raise Exception("Conversion failed: Output PDF not created or empty")
             
        with open(temp_pdf_path, "rb") as f:
            pdf_bytes = f.read()
            
        return io.BytesIO(pdf_bytes)
    
    except Exception as e:
         print(f"DEBUG: Word to PDF error: {str(e)}")
         raise Exception(f"Word to PDF conversion failed: {str(e)}")
         
    finally:
        # Final cleanup safety
        if doc:
            try: doc.Close(False)
            except: pass
        if word:
            try: word.Quit()
            except: pass
        pythoncom.CoUninitialize()
        safe_remove(temp_docx_path)
        safe_remove(temp_pdf_path)

def excel_to_pdf_windows(excel_bytes):
    """
    Convert Excel to PDF using Excel COM automation
    """
    # Get a temporary file path for the input Excel
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as temp_input:
        temp_input.write(excel_bytes)
        temp_input_path = temp_input.name

    # Prepare output PDF path
    output_pdf_path = temp_input_path.replace(".xlsx", ".pdf")

    excel = None
    wb = None
    try:
        # Import win32com.client inside the function to avoid issues on non-Windows
        import win32com.client
        
        # Initialize Excel application
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        
        # Open the workbook
        wb = excel.Workbooks.Open(temp_input_path)
        
        # 0 corresponds to xlTypePDF
        wb.ExportAsFixedFormat(0, output_pdf_path)
        
        # Read the generated PDF
        if os.path.exists(output_pdf_path):
            with open(output_pdf_path, "rb") as pdf_file:
                pdf_bytes = pdf_file.read()
            return io.BytesIO(pdf_bytes)
        else:
            raise Exception("PDF output file was not created by Excel.")
            
    except Exception as e:
        print(f"Excel COM conversion error: {e}")
        # If COM fails or not on Windows, raise error
        raise e
        
    finally:
        # Cleanup
        if wb:
            try:
                wb.Close(SaveChanges=False)
            except:
                pass
        if excel:
            try:
                excel.Quit()
            except:
                pass
                
        safe_remove(temp_input_path)
        safe_remove(output_pdf_path)


def add_page_numbers(pdf_bytes, options):
    """
    Add page numbers to PDF with flexible positioning and styling.
    options: {
        'position': 'bottom-center' (default), top-left, top-center, top-right, bottom-left, bottom-right
        'margin': 'recommended' (default), small, big
        'first_number': 1 (default)
        'page_mode': 'single' (default), facing
        'cover_page': False (default)
    }
    """
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    position = options.get('position', 'bottom-center')
    margin_type = options.get('margin', 'recommended')
    start_number = int(options.get('first_number', 1))
    page_mode = options.get('page_mode', 'single')
    is_cover = options.get('cover_page', False) == 'true' or options.get('cover_page') is True

    # Define margins (in points, 72 pts = 1 inch)
    # A4 is approx 595 x 842 points
    margin_map = {
        'small': 20,
        'recommended': 40,
        'big': 72
    }
    margin = margin_map.get(margin_type, 40)
    
    font_size = 12
    font_color = (0, 0, 0) # Black

    for i, page in enumerate(doc):
        # Skip cover page if requested
        if is_cover and i == 0:
            continue
            
        page_num = start_number + (i - 1 if is_cover else i)
        text = str(page_num)
        
        # Determine position
        rect = page.rect
        width = rect.width
        height = rect.height
        
        # Calculate x, y coordinates
        # Text alignment adjustment (approximate width of text)
        text_width = fitz.get_text_length(text, fontname="helv", fontsize=font_size)
        text_height = font_size
        
        pos_x = 0
        pos_y = 0
        
        # Facing pages logic: Mirror Left/Right if facing
        # If facing: 
        #   Even pages (0, 2, 4...): Left side is outer, Right side is inner? 
        #   Usually: Page 1 (Right), Page 2 (Left), Page 3 (Right)...
        #   But 'i' is 0-indexed.
        current_pos = position
        
        if page_mode == 'facing':
            # Assuming standard book layout:
            # i=0 (Page 1 in PDF viewer) is usually Recto (Right)
            # i=1 (Page 2) is Verso (Left)
            # But let's check widely used logic.
            # Usually users expect flip-flopping.
            # If position is "Bottom-Right", on Left page it should be "Bottom-Left".
            
            # Simple facing logic:
            # If Right side requested (Top-Right, Bottom-Right):
            #   Page i (Even) -> Right
            #   Page i (Odd) -> Left 
            # Wait, usually page 1 is on the right. 
            # If strictly physical pages:
            # Page 1 (Right): margin on right.
            # Page 2 (Left): margin on left.
            
            is_right_page = (i % 2 == 0) # 0, 2, 4... are on the right in 2-up view normally (if cover is separate)
            # Actually simplest is: 
            # If user selected "Right", they want it on the outside corner.
            # Outside corner: Right for Right Page, Left for Left Page.
            
            if "right" in position.lower():
                current_pos = position if is_right_page else position.replace("right", "left")
            elif "left" in position.lower():
                current_pos = position if not is_right_page else position.replace("left", "right")
        
        # Determine coordinates based on current_pos and margin
        if "top" in current_pos:
            pos_y = margin + text_height
        else: # bottom
            pos_y = height - margin
            
        if "left" in current_pos:
            pos_x = margin
        elif "right" in current_pos:
            pos_x = width - margin - text_width
        else: # center
            pos_x = (width - text_width) / 2
            
        # Draw text
        page.insert_text((pos_x, pos_y), text, fontsize=font_size, fontname="helv", color=font_color)

    output_stream = io.BytesIO()
    doc.save(output_stream)
    doc.close()
    output_stream.seek(0)
    return output_stream


def verify_pdf(stream):
    """
    Verify if a byte stream is a valid PDF
    """
    try:
        stream.seek(0)
        import fitz
        doc = fitz.open(stream=stream, filetype="pdf")
        doc.close()
        stream.seek(0)
        return True
    except:
        return False

def verify_docx(stream):
    """
    Verify if a byte stream is a valid DOCX
    """
    try:
        from docx import Document
        stream.seek(0)
        Document(stream)
        stream.seek(0)
        return True
    except:
        return False

def verify_excel(stream):
    """
    Verify if a byte stream is a valid Excel file
    """
    try:
        import pandas as pd
        stream.seek(0)
        pd.read_excel(stream)
        stream.seek(0)
        return True
    except:
        return False

def protect_pdf(pdf_bytes, password):
    """
    Password protect a PDF using PyMuPDF (fitz)
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    if not password:
        raise ValueError("Password is required for protection")
    
    try:
        # Open the PDF from bytes
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Save to memory with encryption
        output_stream = io.BytesIO()
        # permissions: 0 is no permissions except viewing? 
        # Actually fitz.PDF_PERM_PRINT etc are bits. 0 means very restricted.
        # owner_pw can be same as user_pw if not specified.
        doc.save(
            output_stream,
            owner_pw=password,
            user_pw=password,
            encryption=fitz.PDF_ENCRYPT_AES_256, # Use strong encryption
            garbage=3,
            deflate=True
        )
        doc.close()
        
        output_stream.seek(0)
        return output_stream
        
    except Exception as e:
        raise Exception(f"Protect PDF failed: {str(e)}")


def pdf_to_jpg(pdf_bytes):
    """
    Convert PDF pages to JPG images using PyMuPDF (fitz) for speed and accuracy
    Returns a list of BytesIO objects, each containing a JPG image
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        image_streams = []
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            # Use higher zoom for better quality (300 DPI)
            zoom = 2.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            
            img_stream = io.BytesIO(pix.tobytes("jpg"))
            image_streams.append(img_stream)
            
        pdf_document.close()
        return image_streams
        
    except Exception as e:
        raise Exception(f"PDF to JPG conversion failed: {str(e)}")

def pdf_to_png(pdf_bytes):
    """
    Convert PDF pages to PNG images using PyMuPDF (fitz)
    Returns a list of BytesIO objects, each containing a PNG image
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        image_streams = []
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            # Use higher zoom for better quality
            zoom = 2.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            
            img_stream = io.BytesIO(pix.tobytes("png"))
            image_streams.append(img_stream)
            
        pdf_document.close()
        return image_streams
        
    except Exception as e:
        raise Exception(f"PDF to PNG conversion failed: {str(e)}")

def jpg_to_pdf(image_bytes_list):
    """
    Convert multiple JPG images into a single PDF using Pillow
    """
    if not image_bytes_list:
        raise ValueError("No images provided")
    
    try:
        images = []
        for img_bytes in image_bytes_list:
            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            images.append(img)
            
        if not images:
            raise ValueError("Could not process images")
            
        pdf_stream = io.BytesIO()
        images[0].save(pdf_stream, format="PDF", save_all=True, append_images=images[1:], optimize=True)
        pdf_stream.seek(0)
        return pdf_stream
        
    except Exception as e:
        raise Exception(f"JPG to PDF conversion failed: {str(e)}")

def office_to_pdf_linux(file_bytes, input_suffix):
    """
    Convert Office documents to PDF using LibreOffice on Linux
    """
    temp_dir = tempfile.mkdtemp()
    temp_input_path = os.path.join(temp_dir, f"input{input_suffix}")
    
    try:
        with open(temp_input_path, "wb") as f:
            f.write(file_bytes)
            
        office_cmd = 'libreoffice'
        if shutil.which('soffice'):
            office_cmd = 'soffice'
        elif shutil.which('lowriter'):
            office_cmd = 'lowriter'

        cmd = [
            office_cmd, 
            '--headless', 
            '--convert-to', 
            'pdf', 
            '--outdir', 
            temp_dir, 
            temp_input_path
        ]
        
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        
        temp_pdf_path = temp_input_path.rsplit('.', 1)[0] + ".pdf"
        
        if not os.path.exists(temp_pdf_path):
             error_msg = result.stderr.decode() if result.stderr else "Unknown error"
             raise Exception(f"LibreOffice conversion failed. Error: {error_msg}")
            
        with open(temp_pdf_path, "rb") as f:
            pdf_bytes = f.read()
            
        return io.BytesIO(pdf_bytes)
        
    except Exception as e:
        raise Exception(f"Office to PDF conversion failed: {str(e)}")
    finally:
        if os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except:
                pass

def ppt_to_pdf(ppt_bytes):
    """
    Convert PowerPoint to PDF
    """
    system = platform.system().lower()
    if system == 'windows':
        return ppt_to_pdf_windows(ppt_bytes)
    else:
        return office_to_pdf_linux(ppt_bytes, ".pptx")

def safe_remove(file_path, retries=5, delay=0.5):
    """Attempt to remove a file with retries to handle COM lazy handle release"""
    if not file_path or not os.path.exists(file_path):
        return
    for i in range(retries):
        try:
            os.remove(file_path)
            return
        except OSError:
            time.sleep(delay)

def ppt_to_pdf_windows(ppt_bytes):
    """
    Convert PowerPoint to PDF using COM automation for 100% accuracy
    """
    import pythoncom
    import win32com.client
    
    pythoncom.CoInitialize()
    temp_ppt_path = None
    temp_pdf_path = None
    powerpoint = None
    deck = None
    
    try:
        suffix = ".pptx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_ppt:
            temp_ppt.write(ppt_bytes)
            temp_ppt_path = temp_ppt.name
            
        temp_pdf_path = temp_ppt_path.replace(suffix, ".pdf")
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # ppFixedFormatTypePDF = 32
        deck = powerpoint.Presentations.Open(temp_ppt_path, WithWindow=False)
        deck.SaveAs(temp_pdf_path, 32)
        deck.Close()
        powerpoint.Quit()
        
        # Explicitly release COM objects
        del deck
        del powerpoint
        
        with open(temp_pdf_path, "rb") as f:
            pdf_bytes = f.read()
            
        return io.BytesIO(pdf_bytes)
    except Exception as e:
        raise Exception(f"PPT to PDF conversion failed: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        safe_remove(temp_ppt_path)
        safe_remove(temp_pdf_path)

def excel_to_pdf_weasyprint(excel_bytes):
    """
    Fallback Excel to PDF using WeasyPrint. 
    Matches the user's visual request (Grey header, beige rows, "Page1_Table" title).
    """
    import io
    import pandas as pd
    try:
        from weasyprint import HTML, CSS
        print("DEBUG: Using WeasyPrint for Excel to PDF...")
        
        excel_file = io.BytesIO(excel_bytes)
        xl = pd.ExcelFile(excel_file)
        
        html_parts = ["<html><head><style>"]
        html_parts.append("""
            body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; padding: 20px; }
            .page-break { page-break-after: always; }
            .table-title { font-size: 24px; font-weight: bold; margin-bottom: 20px; }
            table { width: 100%; border-collapse: collapse; margin-bottom: 30px; border: 1px solid black; }
            th { background-color: #808080; color: white; padding: 10px; border: 1px solid black; font-weight: bold; }
            td { background-color: #FDFDF0; color: black; padding: 10px; border: 1px solid black; text-align: center; }
        """)
        html_parts.append("</style></head><body>")
        
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            html_parts.append(f'<div class="table-title">{sheet_name if "Table" in sheet_name else f"{sheet_name}_Table"}</div>')
            html_parts.append(df.to_html(index=False, classes='excel-table'))
            html_parts.append('<div class="page-break"></div>')
            
        html_parts.append("</body></html>")
        
        html_content = "".join(html_parts)
        pdf_stream = io.BytesIO()
        HTML(string=html_content).write_pdf(pdf_stream)
        pdf_stream.seek(0)
        return pdf_stream
    except Exception as e:
        print(f"DEBUG: WeasyPrint fallback failed: {e}")
        raise e

def excel_to_pdf(excel_bytes):
    """
    Convert Excel to PDF with styling to match user screenshot.
    Uses win32com on Windows for maximum fidelity, fallback to WeasyPrint.
    Ensures long words are wrapped and not clipped.
    """
    if not excel_bytes:
        raise ValueError("Excel file is empty")
        
    import platform
    system = platform.system().lower()
    
    if system == 'windows':
        import pythoncom
        import win32com.client
        import tempfile
        
        pythoncom.CoInitialize()
        temp_excel = None
        temp_pdf = None
        excel_app = None
        
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as f:
                f.write(excel_bytes)
                temp_excel = f.name
                
            temp_pdf = temp_excel.replace(".xlsx", ".pdf")
            
            # Initialize Excel with caution
            excel_app = com_retry(win32com.client.DispatchEx, "Excel.Application")
            excel_app.Visible = False
            excel_app.DisplayAlerts = False
            excel_app.ScreenUpdating = False
            
            wb = com_retry(excel_app.Workbooks.Open, temp_excel)
            
            # STYLING to match User Screenshot
            sheets = com_retry(getattr, wb, "Sheets")
            for i in range(1, sheets.Count + 1):
                sheet = com_retry(sheets.Item, i)
                
                # 1. Insert Title rows
                sheet_rows = com_retry(getattr, sheet, "Rows")
                com_retry(sheet_rows(1).Insert)
                com_retry(sheet_rows(1).Insert)
                
                sheet_cells = com_retry(getattr, sheet, "Cells")
                title_cell = com_retry(sheet_cells, 1, 1)
                com_retry(setattr, title_cell, "Value", f"{sheet.Name}_Table")
                
                title_font = com_retry(getattr, title_cell, "Font")
                com_retry(setattr, title_font, "Bold", True)
                com_retry(setattr, title_font, "Size", 16)
                
                # 2. Format used range
                used_range = com_retry(getattr, sheet, "UsedRange")
                com_retry(setattr, used_range, "WrapText", True)
                com_retry(setattr, used_range, "VerticalAlignment", -4108) # xlCenter
                
                # Colors/Borders
                header_row = com_retry(sheet_rows, 3)
                header_interior = com_retry(getattr, header_row, "Interior")
                com_retry(setattr, header_interior, "Color", 8421504) # Grey
                
                header_font = com_retry(getattr, header_row, "Font")
                com_retry(setattr, header_font, "Color", 16777215)    # White
                com_retry(setattr, header_font, "Bold", True)
                
                used_rows_count = com_retry(getattr, com_retry(getattr, used_range, "Rows"), "Count")
                used_cols_count = com_retry(getattr, com_retry(getattr, used_range, "Columns"), "Count")
                
                if used_rows_count > 3:
                    data_r1 = com_retry(sheet_cells, 4, 1)
                    data_r2 = com_retry(sheet_cells, used_rows_count, used_cols_count)
                    data_range = com_retry(sheet.Range, data_r1, data_r2)
                    data_interior = com_retry(getattr, data_range, "Interior")
                    com_retry(setattr, data_interior, "Color", 15793661) # Beige
                
                range_borders = com_retry(getattr, used_range, "Borders")
                com_retry(setattr, range_borders, "LineStyle", 1) # xlContinuous
                com_retry(setattr, range_borders, "Weight", 2)    # xlThin
                
                # Layout & Scaling (CRITICAL FOR CLIPPING)
                sheet_cols = com_retry(getattr, sheet, "Columns")
                com_retry(sheet_cols.AutoFit)
                
                for j in range(1, used_cols_count + 1):
                    col = com_retry(sheet_cols.Item, j)
                    curr_width = com_retry(getattr, col, "ColumnWidth")
                    # EXTRA PADDING for long words like "Mathematics"
                    com_retry(setattr, col, "ColumnWidth", curr_width * 1.45) 
                    if col.ColumnWidth > 70: com_retry(setattr, col, "ColumnWidth", 70)
                
                com_retry(sheet_rows.AutoFit)
                
                # PageSetup
                page_setup = com_retry(getattr, sheet, "PageSetup")
                com_retry(setattr, page_setup, "Zoom", False)
                com_retry(setattr, page_setup, "FitToPagesWide", 1)
                com_retry(setattr, page_setup, "FitToPagesTall", False)
                com_retry(setattr, page_setup, "Orientation", 1) # xlPortrait
                com_retry(setattr, page_setup, "CenterHorizontally", True)
                com_retry(setattr, page_setup, "LeftMargin", excel_app.InchesToPoints(0.25))
                com_retry(setattr, page_setup, "RightMargin", excel_app.InchesToPoints(0.25))
            
            com_retry(setattr, excel_app, "ScreenUpdating", True)
            
            # Export
            com_retry(wb.ExportAsFixedFormat, 0, temp_pdf)
            com_retry(wb.Close, False)
            
            if not os.path.exists(temp_pdf) or os.path.getsize(temp_pdf) == 0:
                raise Exception("Excel to PDF conversion failed: PDF not generated or empty")
                
            with open(temp_pdf, "rb") as f:
                pdf_bytes = f.read()
                
            return io.BytesIO(pdf_bytes)
            
        except Exception as e:
            print(f"DEBUG: excel_to_pdf CRASHED: {str(e)}")
            # Fallback to WeasyPrint if win32com fails
            try:
                return excel_to_pdf_weasyprint(excel_bytes)
            except:
                raise Exception(f"Excel to PDF conversion failed: {str(e)}")
        finally:
            if excel_app:
                try: excel_app.Quit()
                except: pass
            safe_remove(temp_excel)
            safe_remove(temp_pdf)
            pythoncom.CoUninitialize()
    else:
        return excel_to_pdf_weasyprint(excel_bytes)


def com_retry(func, *args, **kwargs):
    """Specific wrapper to handle 'Call was rejected by callee' COM errors"""
    import time
    import pywintypes
    last_err = None
    for i in range(30): # Try for 30 seconds
        try:
            return func(*args, **kwargs)
        except pywintypes.com_error as e:
            last_err = e
            hr = getattr(e, 'hresult', None)
            if hr is None and len(e.args) > 0: hr = e.args[0]
            if hr == -2147418111:
                print(f"DEBUG: COM busy (rejected). Retrying {i+1}/30...")
                time.sleep(1)
                continue
            raise e
        except Exception as e:
            last_err = e
            if "-2147418111" in str(e) or "rejected" in str(e).lower():
                print(f"DEBUG: COM busy (string). Retrying {i+1}/30...")
                time.sleep(1)
                continue
            raise e
    if last_err: raise last_err
    return func(*args, **kwargs)

def rotate_pdf(pdf_bytes, angle=90):
    """
    Rotate PDF pages by a specific angle (90, 180, 270)
    This is highly accurate as it only modifies PDF metadata for rotation
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
        
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            # Get current rotation and add the new angle
            current_rotation = page.rotation
            new_rotation = (current_rotation + angle) % 360
            page.set_rotation(new_rotation)
            
        # Save to memory with optimizations
        output_stream = io.BytesIO()
        pdf_document.save(output_stream, garbage=3, deflate=True)
        pdf_document.close()
        
        output_stream.seek(0)
        return output_stream
        
    except Exception as e:
        raise Exception(f"Rotate PDF failed: {str(e)}")

def get_pdf_page_images(pdf_bytes):
    """
    Extract thumbnail images of all PDF pages for preview
    Returns list of base64-encoded JPEG images
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    
    try:
        import base64
        
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        page_images = []
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            
            # Render page at lower resolution for thumbnails (150 DPI)
            zoom = 1.0  # 72 DPI * 1.0 = 72 DPI for thumbnails
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            
            # Convert to JPEG bytes
            img_bytes = pix.tobytes("jpeg")
            
            # Base64 encode for JSON transport
            img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            
            page_images.append({
                'page_number': page_num + 1,
                'image': f"data:image/jpeg;base64,{img_base64}",
                'width': pix.width,
                'height': pix.height
            })
        
        pdf_document.close()
        
        print(f"DEBUG: Generated {len(page_images)} page thumbnails")
        return page_images
        
    except Exception as e:
        raise Exception(f"Get PDF page images failed: {str(e)}")

def delete_pdf_pages(pdf_bytes, pages_to_delete):
    """
    Delete specified pages from PDF
    pages_to_delete: list of page numbers (1-indexed) to remove
    Returns modified PDF as BytesIO stream
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    
    if not pages_to_delete or len(pages_to_delete) == 0:
        raise ValueError("No pages specified for deletion")
    
    try:
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = pdf_document.page_count
        
        # Convert to 0-indexed and sort in reverse order
        pages_to_delete_0indexed = sorted([p - 1 for p in pages_to_delete], reverse=True)
        
        # Validate page numbers
        for page_num in pages_to_delete_0indexed:
            if page_num < 0 or page_num >= total_pages:
                raise ValueError(f"Invalid page number: {page_num + 1}")
        
        # Check if trying to delete all pages
        if len(pages_to_delete) >= total_pages:
            raise ValueError("Cannot delete all pages from PDF")
        
        # Delete pages in reverse order to maintain correct indices
        for page_num in pages_to_delete_0indexed:
            pdf_document.delete_page(page_num)
        
        # Save to memory
        output_stream = io.BytesIO()
        pdf_document.save(output_stream, garbage=3, deflate=True)
        pdf_document.close()
        
        output_stream.seek(0)
        
        print(f"DEBUG: Deleted {len(pages_to_delete)} pages from PDF")
        return output_stream
        
    except Exception as e:
        raise Exception(f"Delete PDF pages failed: {str(e)}")


def compress_pdf(pdf_bytes, level="recommended"):
    """
    Compress PDF using PyMuPDF (fitz) with image downsampling and garbage collection.
    Levels: 'extreme', 'recommended', 'less'
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
        
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Settings based on level
        # Default: Recommended (150 DPI, 75 quality)
        dpi = 150
        quality = 75
        
        if level == "extreme":
            dpi = 72
            quality = 50
        elif level == "less":
            dpi = 200 # Higher DPI
            quality = 90
            
        # Optimize images
        for page in doc:
            for img in page.get_images():
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                
                # If image is larger than target DPI relative to page size? 
                # Simplified: check pixel dimensions vs simple threshold or just blindly re-compress large images
                if pix.width > 600 or pix.height > 600:
                    try:
                        # Create new pixmap with reduced size/quality if needed
                        # Ideally we resample. PyMuPDF 'pix.shrink' or similar is option.
                        # We'll use PIL for better control over efficient JPEG re-encoding
                        
                        # Existing logic was empty. Let's do a robust approach:
                        # 1. Extract image bytes
                        if pix.n - pix.alpha < 4:       # GRAY or RGB
                            img_bytes = pix.tobytes("ppm")
                        else:                           # PMAP: CMYK
                            img_bytes = pix.tobytes("pam")
                            
                        # 2. Open in PIL
                        from PIL import Image
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        
                        # 3. Resize if too big
                        # Calculate target max dimension based on DPI (approx)
                        # Let's say max dimension 1000px for extreme, 2000 for recommended
                        max_dim = 2000
                        if level == "extreme": max_dim = 1000
                        elif level == "less": max_dim = 3000
                        
                        if pil_img.width > max_dim or pil_img.height > max_dim:
                            pil_img.thumbnail((max_dim, max_dim), Image.Resampling.LANCZOS)
                            
                        # 4. Save to buffer as JPEG
                        new_img_stream = io.BytesIO()
                        pil_img.save(new_img_stream, format="JPEG", quality=quality, optimize=True)
                        new_img_stream.seek(0)
                        
                        # 5. Update (replace) image in PDF
                        doc.update_stream(xref, new_img_stream.read())
                        
                    except Exception as img_err:
                        print(f"DEBUG: Image compression failed for xref {xref}: {img_err}")
                        pass

        output_stream = io.BytesIO()
        # garbage=4: dedup images, fonts, etc. deflate=True: compress streams
        doc.save(output_stream, garbage=4, deflate=True, clean=True)
        doc.close()
        
        output_stream.seek(0)
        
        # Check if size actually reduced
        original_size = len(pdf_bytes)
        compressed_size = output_stream.getbuffer().nbytes
        print(f"DEBUG: Compression ({level}). Original: {original_size}, New: {compressed_size}")
        
        if compressed_size >= original_size and level != "less":
             # If fitz didn't help, maybe it's already compressed. 
             # Return original to ensure no quality loss for no gain.
             print("DEBUG: Compressed size larger or equal. Returning original.")
             return io.BytesIO(pdf_bytes)

        return output_stream
        
    except Exception as e:
        raise Exception(f"Compress PDF failed: {str(e)}")

def merge_pdf(pdf_bytes_list):
    """
    Merge multiple PDF files into one using PyMuPDF (fitz) which is more robust
    """
    if not pdf_bytes_list:
        raise ValueError("No PDF files provided")
        
    try:
        merged_doc = fitz.open()
        
        for pdf_bytes in pdf_bytes_list:
             with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                 merged_doc.insert_pdf(doc)
            
        output_stream = io.BytesIO()
        merged_doc.save(output_stream, garbage=3, deflate=True)
        merged_doc.close()
            
        output_stream.seek(0)
        return output_stream
        
    except Exception as e:
        raise Exception(f"Merge PDF failed: {str(e)}")

def pdf_to_ppt(pdf_bytes):
    """
    Convert PDF to PowerPoint (PPTX) by converting pages to images and placing on slides.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import io
        
        # 1. Convert PDF to List of Images
        images = pdf_to_jpg(pdf_bytes) # Reusing existing function
        
        if not images:
             raise Exception("Failed to extract images from PDF")

        # 2. Create PPTX
        prs = Presentation()
        # Use a blank slide layout (Index 6 is usually blank)
        blank_slide_layout = prs.slide_layouts[6]
        
        for img_stream in images:
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Use temporary file for image because python-pptx works best with file paths
            # (It can handle blobs but temp file handles types better sometimes)
            img_stream.seek(0)
            
            # Determine image size to maintain aspect ratio or fit slide?
            # We will fit slide.
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            slide.shapes.add_picture(img_stream, 0, 0, width=slide_width, height=slide_height)

        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return output_stream

    except Exception as e:
        raise Exception(f"PDF to PPT conversion failed: {str(e)}")

def excel_to_pdf_windows(excel_bytes):
    """
    Convert Excel to PDF using Excel COM automation
    """
    import pythoncom
    import win32com.client
    
    pythoncom.CoInitialize()
    temp_xl_path = None
    temp_pdf_path = None
    excel = None
    wb = None
    
    try:
        suffix = ".xlsx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_xl:
            temp_xl.write(excel_bytes)
            temp_xl_path = temp_xl.name
        
        # Convert to absolute path
        temp_xl_path = os.path.abspath(temp_xl_path)
        temp_pdf_path = temp_xl_path.replace(suffix, ".pdf")
        
        print(f"DEBUG: Excel temp file: {temp_xl_path}")
        print(f"DEBUG: PDF output file: {temp_pdf_path}")
        
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        
        print(f"DEBUG: Opening Excel file...")
        wb = excel.Workbooks.Open(temp_xl_path)
        
        print(f"DEBUG: Exporting to PDF...")
        # xlTypePDF = 0
        wb.ExportAsFixedFormat(0, temp_pdf_path)
        
        print(f"DEBUG: Closing workbook...")
        wb.Close(SaveChanges=False)
        excel.Quit()
        
        del wb
        del excel
        
        print(f"DEBUG: Reading PDF file...")
        if not os.path.exists(temp_pdf_path):
            raise Exception("PDF file was not created")
            
        with open(temp_pdf_path, "rb") as f:
            pdf_bytes = f.read()
        
        print(f"DEBUG: PDF size: {len(pdf_bytes)} bytes")
        return io.BytesIO(pdf_bytes)
        
    except Exception as e:
        print(f"DEBUG: Excel to PDF error: {str(e)}")
        raise Exception(f"Excel to PDF failed: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        safe_remove(temp_xl_path)
        safe_remove(temp_pdf_path)




def excel_to_pdf_fallback(excel_bytes):
    """
    Fallback Excel to PDF converter using openpyxl and reportlab
    Used when COM automation is not available
    """
    try:
        from openpyxl import load_workbook
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        
        print("DEBUG: Using fallback Excel to PDF converter...")
        
        # Load Excel workbook
        wb = load_workbook(io.BytesIO(excel_bytes), data_only=True)
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        elements = []
        styles = getSampleStyleSheet()
        
        # Process each sheet
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            
            # Add sheet title
            if sheet_idx > 0:
                elements.append(PageBreak())
            
            title = Paragraph(f"<b>{sheet_name}</b>", styles['Heading1'])
            elements.append(title)
            elements.append(Spacer(1, 0.2*inch))
            
            # Extract data from sheet
            data = []
            for row in ws.iter_rows(values_only=True):
                # Convert None to empty string and handle other types
                row_data = []
                for cell in row:
                    if cell is None:
                        row_data.append('')
                    else:
                        row_data.append(str(cell))
                data.append(row_data)
            
            if data:
                # Create table
                table = Table(data)
                
                # Style the table
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 1), (-1, -1), 8),
                ]))
                
                elements.append(table)
        
        # Build PDF
        doc.build(elements)
        pdf_buffer.seek(0)
        
        print("DEBUG: Fallback Excel to PDF conversion successful")
        return pdf_buffer
        
    except Exception as e:
        print(f"DEBUG: Fallback Excel to PDF error: {str(e)}")
        raise Exception(f"Fallback Excel to PDF conversion failed: {str(e)}")


def excel_to_pdf(excel_bytes):
    """
    Convert Excel to PDF using platform-specific methods
    - Windows: Tries Excel COM automation first, falls back to openpyxl+reportlab
    - Linux: Uses LibreOffice for compatibility
    """
    print(f"DEBUG: Starting Excel to PDF conversion... Size: {len(excel_bytes)} bytes")
    system = platform.system().lower()
    
    if system == 'windows':
        try:
            result = excel_to_pdf_windows(excel_bytes)
            print("DEBUG: Excel to PDF conversion complete (COM).")
            return result
        except Exception as com_error:
            print(f"DEBUG: COM automation failed: {str(com_error)}")
            print("DEBUG: Trying fallback method...")
            try:
                result = excel_to_pdf_fallback(excel_bytes)
                print("DEBUG: Excel to PDF conversion complete (fallback).")
                return result
            except Exception as fallback_error:
                print(f"DEBUG: Fallback also failed: {str(fallback_error)}")
                raise Exception(f"Excel to PDF conversion failed. COM error: {str(com_error)}, Fallback error: {str(fallback_error)}")
    else:
        result = office_to_pdf_linux(excel_bytes, ".xlsx")
        print("DEBUG: Excel to PDF conversion complete (LibreOffice).")
        return result



def is_java_available():
    """Check if Java is available for tabula-py"""
    import subprocess
    try:
        subprocess.run(['java', '-version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        return True
    except:
        return False

def pdf_to_excel(pdf_bytes):
    """
    Convert PDF to Excel with extreme robustness and visual styling.
    Ensures a valid file is ALWAYS returned.
    """
    if not pdf_bytes:
        raise ValueError("PDF file is empty")
    
    import io
    import pandas as pd
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Border, Side, Alignment
    
    wb = Workbook()
    wb.remove(wb.active)
    
    any_content = False
    java_ok = is_java_available()
    
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                tables = []
                
                # 1. Try Tabula
                if java_ok:
                    try:
                        import tabula
                        dfs = tabula.read_pdf(io.BytesIO(pdf_bytes), pages=i, multiple_tables=True, silent=True)
                        if dfs:
                            for df in dfs:
                                if not df.empty: tables.append(df)
                    except: pass
                
                # 2. Try pdfplumber
                try:
                    pqt = page.extract_tables()
                    if pqt:
                        for t in pqt:
                            if t:
                                df = pd.DataFrame(t)
                                df = df.dropna(how='all').dropna(axis=1, how='all')
                                if not df.empty: tables.append(df)
                except: pass
                
                # 3. Advanced Text fallback (Word Clustering)
                if not tables:
                    try:
                        words = page.extract_words()
                        if words:
                            # Cluster words into lines based on Y coordinate
                            lines = {}
                            for w in words:
                                y = round(float(w['top']), 1)
                                if y not in lines: lines[y] = []
                                lines[y].append(w)
                            
                            row_data = []
                            for y in sorted(lines.keys()):
                                line_words = sorted(lines[y], key=lambda x: float(x['x0']))
                                # Cluster these words into columns based on X gap
                                current_row = []
                                if line_words:
                                    current_part = line_words[0]['text']
                                    for idx in range(1, len(line_words)):
                                        gap = float(line_words[idx]['x0']) - float(line_words[idx-1]['x1'])
                                        if gap > 3: # Threshold for new column (Lowered for precision)
                                            current_row.append(current_part)
                                            current_part = line_words[idx]['text']
                                        else:
                                            current_part += " " + line_words[idx]['text']
                                    current_row.append(current_part)
                                if current_row:
                                    row_data.append(current_row)
                            
                            if row_data:
                                # Standardize row lengths for DataFrame
                                max_cols = max(len(r) for r in row_data)
                                padded_rows = [r + [""] * (max_cols - len(r)) for r in row_data]
                                # Use first row as header if it exists
                                df = pd.DataFrame(padded_rows)
                                if not df.empty and len(df) > 1:
                                    df.columns = df.iloc[0]
                                    df = df[1:]
                                tables.append(df)
                    except Exception as fallback_err:
                        print(f"DEBUG: Advanced fallback failed: {fallback_err}")
                
                ws = wb.create_sheet(title=f"Sheet{i}")
                ws.cell(row=1, column=1, value=f"Sheet{i}").font = Font(bold=True, size=16)
                
                curr_row = 3
                if tables:
                    any_content = True
                    for df in tables:
                        # Header
                        for c_idx, val in enumerate(df.columns, 1):
                            cell = ws.cell(row=curr_row, column=c_idx, value=str(val))
                            cell.fill = PatternFill(start_color="808080", end_color="808080", fill_type="solid")
                            cell.font = Font(color="FFFFFF", bold=True)
                            cell.alignment = Alignment(horizontal='center', wrap_text=True)
                        curr_row += 1
                        # Data
                        for _, row in df.iterrows():
                            for c_idx, val in enumerate(row, 1):
                                cell = ws.cell(row=curr_row, column=c_idx, value=str(val) if val is not None else "")
                                cell.fill = PatternFill(start_color="FDFDF0", end_color="FDFDF0", fill_type="solid")
                                cell.alignment = Alignment(horizontal='center', wrap_text=True)
                                cell.border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
                            curr_row += 1
                        curr_row += 2
                else:
                    ws.cell(row=3, column=1, value="No tabular content detected.")

        if not any_content:
            ws = wb.create_sheet(title="No Content")
            ws.cell(row=1, column=1, value="This PDF appears to be an image or contains no extractable data.")

        # Autofit & Styling
        for sheet in wb.worksheets:
            for col in sheet.columns:
                col_letter = col[0].column_letter
                sheet.column_dimensions[col_letter].width = 25 # Default robust width
            sheet.sheet_view.showGridLines = False

        stream = io.BytesIO()
        wb.save(stream)
        stream.seek(0)
        return stream

    except Exception as e:
        # Fallback to absolute minimum valid workbook
        print(f"FATAL PDF_TO_EXCEL: {e}")
        err_wb = Workbook()
        err_wb.active.cell(row=1, column=1, value="Conversion Error: " + str(e))
        s = io.BytesIO()
        err_wb.save(s)
        s.seek(0)
        return s


def verify_excel(stream):
    """Verify that the Excel stream is valid."""
    try:
        if not stream: return False
        stream.seek(0)
        # Check PK signature (XLSX)
        header = stream.read(4)
        stream.seek(0)
        return header == b'PK\x03\x04'
    except:
        return False


def verify_docx(docx_stream):
    """Verify that the DOCX stream contains a valid Word document"""
    try:
        docx_stream.seek(0)
        Document(docx_stream)
        docx_stream.seek(0)
        return True
    except Exception as e:
        print(f"DOCX Validation Error: {e}")
        return False

def verify_pdf(pdf_stream):
    """Verify that the PDF stream contains a valid PDF document"""
    try:
        pdf_stream.seek(0)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        valid = doc.page_count > 0
        doc.close()
        pdf_stream.seek(0)
        return valid
    except Exception as e:
        print(f"PDF Validation Error: {e}")
        return False

